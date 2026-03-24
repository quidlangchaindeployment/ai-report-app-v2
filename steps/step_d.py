# steps/step_d.py
# --- Step D: PowerPoint 出力（v2） -------------------------------------------------
# 役割：
# - Step C で作成したレポートJSON（章/スライド構成）を入力
# - （任意）テンプレPPTXをアップロードして、指定レイアウトに割り当て
# - python-pptx でスライドを自動生成し、ダウンロード可能に
# - （任意）LLM（services.llm.get_llm）があれば、見出し/箇条書きの言い回しを軽く整える
# -----------------------------------------------------------------------------

from __future__ import annotations

import io
import json
from typing import Any, Dict, List, Optional, Callable

import streamlit as st

# python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

# LLM（任意）
_HAS_LLM = False
try:
    from services.llm import get_llm  # Gemini クライアント
    _HAS_LLM = True
except Exception:
    get_llm = None


# --- 1) LLM での軽い言い回し整形（任意） ------------------------------------------
def _ai_polish(text: str, tone_hint: str = "") -> str:
    """
    見出しや箇条書きの言い回しを軽く整える。LLM未設定の場合は原文を返す。
    """
    if not text.strip():
        return text
    if not _HAS_LLM or get_llm is None:
        return text

    llm = get_llm(model_name="gemini-2.5-flash-lite", temperature=0.1, timeout_seconds=30)
    if llm is None:
        return text

    prompt = (
        "次の文を箇条書きのトーンで、与えられた文体に軽く整えてください。意味は変えないでください。\n"
        f"【トーンの指示】{tone_hint or 'ビジネス/経営会議向け・簡潔'}\n"
        f"---\n{text}\n---\n"
        "出力はテキストのみ。"
    )
    try:
        try:
            out = llm.invoke(prompt)  # type: ignore
        except Exception:
            out = llm.predict(prompt)  # type: ignore
        return str(out).strip()
    except Exception:
        return text


# --- 2) レポートJSONのバリデーション ----------------------------------------------
def _validate_report_json(obj: Dict[str, Any]) -> Optional[str]:
    if not isinstance(obj, dict):
        return "JSONの最上位がオブジェクトではありません。"
    if "title" not in obj or "slides" not in obj:
        return "JSON に 'title' または 'slides' がありません。"
    if not isinstance(obj["slides"], list) or not obj["slides"]:
        return "'slides' が空、または配列ではありません。"
    return None


# --- 3) スライド生成のユーティリティ ----------------------------------------------
def _add_title_slide(prs: Presentation, heading: str, bullets: Optional[List[str]] = None):
    # レイアウト index はテンプレにより異なるため、Title (0) を優先
    layout = prs.slide_layouts[0 if len(prs.slide_layouts) > 0 else 0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = heading[:120]
    # サブタイトル/本文枠（存在すれば）
    body = None
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False) and shp != slide.shapes.title:
            body = shp.text_frame
            break
    if body and bullets:
        body.clear()
        first = True
        for b in bullets[:5]:
            if first:
                body.text = b[:200]
                first = False
            else:
                p = body.add_paragraph()
                p.text = b[:200]


def _add_toc_slide(prs: Presentation, heading: str, entries: List[str]):
    # Title and Content(1) を想定
    layout = prs.slide_layouts[1 if len(prs.slide_layouts) > 1 else 0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = heading[:120]
    body = None
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False) and shp != slide.shapes.title:
            body = shp.text_frame
            break
    if body:
        body.clear()
        first = True
        for e in entries[:10]:
            if first:
                body.text = f"• {e[:200]}"
                first = False
            else:
                p = body.add_paragraph()
                p.text = f"• {e[:200]}"


def _add_bullets_slide(prs: Presentation, heading: str, bullets: List[str]):
    layout = prs.slide_layouts[1 if len(prs.slide_layouts) > 1 else 0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = heading[:120]
    body = None
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False) and shp != slide.shapes.title:
            body = shp.text_frame
            break
    if body:
        body.clear()
        first = True
        for b in bullets[:12]:
            if first:
                body.text = f"• {b[:250]}"
                first = False
            else:
                p = body.add_paragraph()
                p.text = f"• {b[:250]}"


def _add_table_slide(prs: Presentation, heading: str, table_dict: Dict[str, Any]):
    # Title Only(5) を想定、無ければ Title and Content(1)
    layout = prs.slide_layouts[5 if len(prs.slide_layouts) > 5 else (1 if len(prs.slide_layouts) > 1 else 0)]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = heading[:120]

    # {"colA": ["v1","v2"], "colB": [...]} を想定
    columns = list(table_dict.keys())
    rows = max([len(v) for v in table_dict.values() if isinstance(v, list)] or [0])
    if not columns or rows == 0:
        # フォールバック：本文にテキストで提示
        body = None
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and shp != slide.shapes.title:
                body = shp.text_frame
                break
        if body:
            body.clear()
            first = True
            for c in columns[:10]:
                val = table_dict.get(c)
                if isinstance(val, list):
                    txt = f"{c}: {', '.join(map(str, val[:5]))}"
                else:
                    txt = f"{c}: {str(val)[:100]}"
                if first:
                    body.text = txt
                    first = False
                else:
                    p = body.add_paragraph()
                    p.text = txt
        return

    # 表の挿入
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(5.0)
    try:
        tbl_shape = slide.shapes.add_table(rows + 1, len(columns), left, top, width, height)
        tbl = tbl_shape.table
        # ヘッダ
        for j, col_name in enumerate(columns):
            tbl.cell(0, j).text = str(col_name)
        # データ
        for i in range(rows):
            for j, col_name in enumerate(columns):
                vals = table_dict.get(col_name, [])
                val = vals[i] if isinstance(vals, list) and i < len(vals) else ""
                tbl.cell(i + 1, j).text = str(val)[:200]
    except Exception:
        # 失敗時は箇条書きにフォールバック
        _add_bullets_slide(prs, heading, ["（表の描画に失敗したため、箇条書きにフォールバックしました）"])


# --- 4) JSON → スライド構築（polish_fn を引数で受け取る） ------------------------
def _build_presentation_from_report(
    report_json: Dict[str, Any],
    template_bytes: Optional[bytes],
    polish_fn: Optional[Callable[[str], str]] = None,
) -> bytes:
    """
    report_json: Step C で生成した {"title": "...", "slides": [...]}
    template_bytes: .pptx のバイナリ（None の場合は白紙テンプレ）
    polish_fn: 見出し・箇条書きの軽い整形を行う関数（None なら無加工）
    """
    if not _HAS_PPTX:
        raise RuntimeError("python-pptx がインストールされていません。requirements.txt を確認してください。")

    prs = Presentation(io.BytesIO(template_bytes)) if (template_bytes is not None) else Presentation()

    title = report_json.get("title", "AI レポート")
    slides = report_json.get("slides", [])
    # デフォルトの整形器
    if polish_fn is None:
        def polish_fn(x: str) -> str:
            return x

    for s in slides:
        layout = str(s.get("layout", "bullets")).lower()
        heading = polish_fn(str(s.get("heading", title))[:120])

        if layout == "title":
            bullets = s.get("bullets", [])
            bullets = [polish_fn(b) for b in bullets] if isinstance(bullets, list) else []
            _add_title_slide(prs, heading, bullets)
        elif layout == "toc":
            toc_items = s.get("bullets", [])
            toc_items = [str(x) for x in toc_items] if isinstance(toc_items, list) else []
            _add_toc_slide(prs, heading, toc_items)
        elif layout == "table":
            table_dict = s.get("table", {})
            if isinstance(table_dict, dict):
                _add_table_slide(prs, heading, table_dict)
            else:
                _add_bullets_slide(prs, heading, ["（table が不正のため箇条書きにフォールバック）"])
        else:
            bullets = s.get("bullets", [])
            bullets = [polish_fn(b) for b in bullets] if isinstance(bullets, list) else []
            _add_bullets_slide(prs, heading, bullets)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# --- 5) UI 本体 -------------------------------------------------------------------
def render():
    st.title("📑 Step D: PowerPoint 出力（v2）")

    if not _HAS_PPTX:
        st.error("python-pptx が見つかりません。`pip install python-pptx` を実行してください。")
        return

    # 1) レポートJSONの入力
    st.header("1) レポートJSONの読み込み")
    uploaded_report = st.file_uploader("report_for_powerpoint.json をアップロード", type=["json"])
    report_json: Optional[Dict[str, Any]] = None

    default_text = ""
    if "step_c_report_json" in st.session_state and st.session_state["step_c_report_json"]:
        default_text = json.dumps(st.session_state["step_c_report_json"], ensure_ascii=False, indent=2)

    text_area = st.text_area("（任意）JSONを直接貼り付け", value=default_text, height=220)

    if uploaded_report:
        try:
            raw = uploaded_report.read().decode("utf-8")
            report_json = json.loads(raw)
            st.success("アップロードしたレポートJSONを使用します。")
        except Exception as e:
            st.error(f"JSON の読み込みに失敗: {e}")
            return
    else:
        if text_area.strip():
            try:
                report_json = json.loads(text_area)
                st.info("貼り付けられたレポートJSONを使用します。")
            except Exception as e:
                st.error(f"JSON の解析に失敗: {e}")
                return
        else:
            st.warning("レポートJSONをアップロードするか、テキスト欄に貼り付けてください。")
            return

    err = _validate_report_json(report_json)
    if err:
        st.error(f"レポートJSONが不正: {err}")
        return

    # 2) テンプレPPTX（任意）
    st.markdown("---")
    st.header("2) テンプレPPTX（任意）を読み込む")
    uploaded_tmpl = st.file_uploader("テンプレート（.pptx）", type=["pptx"], key="tmpl_uploader")
    tmpl_bytes = uploaded_tmpl.read() if uploaded_tmpl else None

    # 3) 生成オプション
    st.markdown("---")
    st.header("3) 生成オプション")
    tone_hint = st.text_input("（任意）AI整形のトーン", value="ビジネス・簡潔・結論先出し")
    use_ai_polish = st.checkbox("見出し/箇条書きをAIで軽く整形する", value=_HAS_LLM)

    # polish_fn を準備
    def _polish_fn_local(s: str) -> str:
        return _ai_polish(s, tone_hint) if use_ai_polish else s

    # 4) 生成
    st.markdown("---")
    st.header("4) PowerPoint を生成")
    if st.button("📤 生成してダウンロード", type="primary"):
        with st.spinner("PowerPoint を生成しています..."):
            ppt_bytes = _build_presentation_from_report(
                report_json=report_json,
                template_bytes=tmpl_bytes,
                polish_fn=_polish_fn_local,
            )
        st.success("PowerPoint の生成に成功しました。下のボタンからダウンロードできます。")
        st.download_button(
            "🔽 PPTX をダウンロード",
            data=ppt_bytes,
            file_name="ai_report_output.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )

    # 5) プレビュー（簡易）
    st.markdown("---")
    st.header("5) レポートJSONのプレビュー")
    with st.expander("JSON（再掲）", expanded=False):
        st.json(report_json)